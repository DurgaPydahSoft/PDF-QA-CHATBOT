import pypdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import io

def extract_text_from_pdf(file_stream) -> str:
    text = ""
    try:
        reader = pypdf.PdfReader(file_stream)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text.strip()

def extract_text_from_docx(file_stream) -> str:
    text = ""
    try:
        doc = Document(file_stream)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    return text.strip()

def extract_text_from_xlsx(file_stream) -> str:
    text = ""
    try:
        wb = load_workbook(file_stream, data_only=True)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text += row_text + "\n"
    except Exception as e:
        print(f"Error extracting text from XLSX: {e}")
    return text.strip()

def extract_text_from_pptx(file_stream) -> str:
    text = ""
    try:
        prs = Presentation(file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    return text.strip()

def extract_text(file_content: bytes, filename: str) -> str:
    ext = filename.split('.')[-1].lower()
    file_stream = io.BytesIO(file_content)
    
    if ext == 'pdf':
        return extract_text_from_pdf(file_stream)
    elif ext == 'docx':
        return extract_text_from_docx(file_stream)
    elif ext in ['xlsx', 'xls']:
        return extract_text_from_xlsx(file_stream)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_stream)
    else:
        print(f"Unsupported file extension: {ext}")
        return ""
