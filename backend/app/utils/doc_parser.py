import pypdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import io

# Try to import PyMuPDF at module level for better error handling
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("Warning: PyMuPDF (fitz) not available. Will use pypdf as fallback for PDF extraction.")

def extract_text_from_pdf(file_stream) -> str:
    text = ""
    
    # Try PyMuPDF first (better quality extraction)
    if PYMUPDF_AVAILABLE:
        try:
            # fitz requires bytes, not a BytesIO object directly
            # Read the stream content as bytes
            file_bytes = file_stream.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            print(f"DEBUG: PDF has {len(doc)} pages.")
            for i, page in enumerate(doc):
                page_text = page.get_text()
                text += page_text + "\n"
                if i % 10 == 0:
                    print(f"DEBUG: Parsed page {i+1}, text len so far: {len(text)}")
            
            print(f"DEBUG: Final extracted text length: {len(text)}")
            doc.close()
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PDF with PyMuPDF: {e}")
            # Fall through to pypdf fallback
    
    # Fallback to pypdf
    try:
        file_stream.seek(0)
        reader = pypdf.PdfReader(file_stream)
        for page in reader.pages:
            text += page.extract_text() + "\n"
        print("DEBUG: Used pypdf fallback for PDF extraction.")
    except Exception as e2:
        print(f"Fallback PDF extraction with pypdf also failed: {e2}")

    return text.strip()

def extract_text_from_docx(file_stream) -> str:
    text = ""
    try:
        doc = Document(file_stream)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    return text.strip()

def extract_text_from_xlsx(file_stream) -> str:
    text = ""
    try:
        wb = load_workbook(file_stream, data_only=True)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text += row_text + "\n"
    except Exception as e:
        print(f"Error extracting text from XLSX: {e}")
    return text.strip()

def extract_text_from_pptx(file_stream) -> str:
    text = ""
    try:
        prs = Presentation(file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    return text.strip()

def extract_text(file_content: bytes, filename: str) -> str:
    import os
    _, ext = os.path.splitext(filename.lower())
    ext = ext.lstrip('.')
    file_stream = io.BytesIO(file_content)
    
    if ext == 'pdf':
        return extract_text_from_pdf(file_stream)
    elif ext == 'docx':
        return extract_text_from_docx(file_stream)
    elif ext in ['xlsx', 'xls']:
        return extract_text_from_xlsx(file_stream)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_stream)
    elif ext == 'txt':
        try:
            return file_content.decode('utf-8')
        except:
            return file_content.decode('latin-1')
    else:
        print(f"Unsupported file extension for '{filename}': {ext}")
        return ""
